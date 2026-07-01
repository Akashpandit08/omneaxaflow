import os
import tempfile
from typing import Dict, Any, List
from pptx import Presentation
import pymupdf  # fitz

class ImportService:
    @staticmethod
    def parse_pptx(file_path: str) -> Dict[str, Any]:
        prs = Presentation(file_path)
        scenes = []
        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
            
            slide_text = "\n".join(text_runs).strip()
            
            scenes.append({
                "id": f"scene-{i+1}",
                "title": f"Slide {i+1}",
                "text": slide_text,
                "script": slide_text,  # Use slide text as basic script
                "duration": 5,
                "transition": "fade",
                "avatar_id": None,
                "voice_id": None
            })
            
        return {
            "title": "Imported PPTX Project",
            "scenes": scenes
        }

    @staticmethod
    def parse_pdf(file_path: str) -> Dict[str, Any]:
        doc = pymupdf.open(file_path)
        scenes = []
        for i, page in enumerate(doc):
            text = page.get_text().strip()
            scenes.append({
                "id": f"scene-{i+1}",
                "title": f"Page {i+1}",
                "text": text,
                "script": text,
                "duration": 5,
                "transition": "fade",
                "avatar_id": None,
                "voice_id": None
            })
            
        return {
            "title": "Imported PDF Project",
            "scenes": scenes
        }

    @staticmethod
    def parse_file(file_path: str, file_type: str) -> Dict[str, Any]:
        if file_type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint", ".pptx", ".ppt"]:
            return ImportService.parse_pptx(file_path)
        elif file_type in ["application/pdf", ".pdf"]:
            return ImportService.parse_pdf(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
